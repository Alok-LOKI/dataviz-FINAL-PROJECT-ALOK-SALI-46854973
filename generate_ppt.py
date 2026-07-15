import collections
import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
# Set widescreen 16:9
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Color Scheme
BG_DARK = RGBColor(9, 9, 11)      # Zinc 950 (#09090b)
BG_LIGHT = RGBColor(248, 250, 252) # Slate 50 (#f8fafc)
TEXT_LIGHT = RGBColor(243, 244, 246) # Gray 100 (#f3f4f6)
TEXT_DARK = RGBColor(15, 23, 42)    # Slate 900 (#0f172a)
ACCENT_AMBER = RGBColor(194, 65, 12) # Orange 700 (#c2410c)
ACCENT_BLUE = RGBColor(2, 132, 199) # Sky 600 (#0284c7)
MUTED_GRAY = RGBColor(148, 163, 184) # Slate 400 (#94a3b8)

# Helper function to set slide solid background
def set_bg_color(slide, rgb_color):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = rgb_color

# Helper to add text block
def add_textbox(slide, left, top, width, height, text, font_size=14, bold=False, color=TEXT_DARK, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = 'Arial'
    return txBox

# Helper to add modern slide titles
def add_slide_header(slide, title_text, dark_mode=False):
    color = TEXT_LIGHT if dark_mode else TEXT_DARK
    add_textbox(slide, Inches(0.8), Inches(0.5), Inches(11.7), Inches(0.8), title_text, font_size=28, bold=True, color=color)

# SLIDE 1: Title Slide (Dark)
slide_layout = prs.slide_layouts[6] # Blank
slide = prs.slides.add_slide(slide_layout)
set_bg_color(slide, BG_DARK)

# Title Text Frame
title_box = slide.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.3), Inches(3.0))
tf = title_box.text_frame
tf.word_wrap = True

p = tf.paragraphs[0]
p.text = "Decarbonization Dynamics"
p.font.size = Pt(48)
p.font.bold = True
p.font.color.rgb = TEXT_LIGHT
p.font.name = 'Arial'

p2 = tf.add_paragraph()
p2.text = "A Multi-Dimensional Story of Global CO2 Emissions"
p2.font.size = Pt(20)
p2.font.color.rgb = MUTED_GRAY
p2.font.name = 'Arial'
p2.space_before = Pt(12)

p3 = tf.add_paragraph()
p3.text = "Final Individual Project • Data Visualization • Summer 2026"
p3.font.size = Pt(14)
p3.font.color.rgb = ACCENT_AMBER
p3.font.name = 'Arial'
p3.space_before = Pt(40)


# SLIDE 2: Project Overview (Light)
slide = prs.slides.add_slide(slide_layout)
set_bg_color(slide, BG_LIGHT)
add_slide_header(slide, "Project Overview & Methodology")

# Left Column (Core Details)
left_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(5.5), Inches(5.0))
tf = left_box.text_frame
tf.word_wrap = True

def add_bullet(tf, title, desc, space=12):
    p = tf.add_paragraph() if tf.paragraphs[0].text else tf.paragraphs[0]
    p.text = title + ": "
    p.font.bold = True
    p.font.size = Pt(15)
    p.font.color.rgb = TEXT_DARK
    p.font.name = 'Arial'
    
    # Add normal text to the same paragraph
    run = p.add_run()
    run.text = desc
    run.font.bold = False
    run.font.size = Pt(15)
    run.font.color.rgb = TEXT_DARK
    
    p.space_after = Pt(space)

add_bullet(tf, "Dataset Sourced", "Our World in Data (OWID) CO2 & GHG Emissions dataset. A rich, high-fidelity real-world dataset covering historical variables globally.")
add_bullet(tf, "Analytical Focus", "Examining shifts in absolute regional emissions, decoupling pathways (GDP vs CO2 per capita), and energy transition structures.")
add_bullet(tf, "Visual Standards", "Plotly-exclusive figures incorporating CVD-safe color mappings, clear takeaway-oriented titles, direct annotations, and decluttered layouts.")

# Right Column (Deliverables summary)
right_box = slide.shapes.add_textbox(Inches(6.8), Inches(1.5), Inches(5.5), Inches(5.0))
tf_right = right_box.text_frame
tf_right.word_wrap = True

add_bullet(tf_right, "Analysis Notebook", "An structured Python notebook containing data loading, exploratory analysis, and 10 detailed publication-ready figures.")
add_bullet(tf_right, "Interactive Dashboard", "A Streamlit web application showcasing a curated interactive set of the findings with dynamic sidebar controls and tabs.")
add_bullet(tf_right, "Core Conclusions", "Evidence shows successful decoupling of emissions and economic growth in developed countries, but rising absolute volumes in developing regions.")


# SLIDE 3: Regional Shifts (Light + Image)
slide = prs.slides.add_slide(slide_layout)
set_bg_color(slide, BG_LIGHT)
add_slide_header(slide, "1. Global Shifts: The Rise of Asian Emissions")

# Text Left
left_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(4.8), Inches(5.0))
tf = left_box.text_frame
tf.word_wrap = True
add_bullet(tf, "Historical Context", "For over a century, North America and Europe produced the vast majority of industrial CO2 emissions.")
add_bullet(tf, "Modern Shift", "Since 1970, Asian emissions have expanded exponentially to support rapid industrialization and population growth.")
add_bullet(tf, "Current State", "Asia now contributes more annual emissions than the rest of the world combined, while Western regions have stabilized and begun a slow decline.")

# Image Right
img_path = r"C:\Users\aloks\.gemini\antigravity-ide\brain\ddd039d2-cd98-4d4a-876a-8a7511b866fd\global_trends_tab_1784127624483.png"
slide.shapes.add_picture(img_path, Inches(5.8), Inches(1.5), width=Inches(6.8))


# SLIDE 4: Economic Decoupling (Light + Image)
slide = prs.slides.add_slide(slide_layout)
set_bg_color(slide, BG_LIGHT)
add_slide_header(slide, "2. Economic Decoupling: G7 vs. BRICS")

# Text Left
left_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(4.8), Inches(5.0))
tf = left_box.text_frame
tf.word_wrap = True
add_bullet(tf, "Western Decoupling", "Major G7 economies (e.g., USA, UK, Germany) show clear 'decoupling' post-2000. GDP per capita grows (dotted lines) while CO2 per capita falls (solid lines).")
add_bullet(tf, "Leapfrog Potential", "Developing nations show rising emissions in tandem with wealth, but present opportunities to transition directly to clean energy systems.")
add_bullet(tf, "Visual Analysis", "Normalizing to base index (2000 = 100) displays this divergence clearly.")

# Image Right
img_path = r"C:\Users\aloks\.gemini\antigravity-ide\brain\ddd039d2-cd98-4d4a-876a-8a7511b866fd\economic_decoupling_tab_1784127635172.png"
slide.shapes.add_picture(img_path, Inches(5.8), Inches(1.5), width=Inches(6.8))


# SLIDE 5: Fuel Sources (Light + Image)
slide = prs.slides.add_slide(slide_layout)
set_bg_color(slide, BG_LIGHT)
add_slide_header(slide, "3. Fuel Breakdown & Resource Transitions")

# Text Left
left_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(4.8), Inches(5.0))
tf = left_box.text_frame
tf.word_wrap = True
add_bullet(tf, "Dominant Resources", "Coal and Oil remain the leading drivers of absolute emissions for the top 10 global emitters.")
add_bullet(tf, "Coal Dependency", "Nations like China and India have a highly coal-dominated emission profile due to heavy reliance on coal power plants.")
add_bullet(tf, "Gas Transition", "Developed nations have significantly reduced coal emissions by switching to natural gas as a bridge fuel, alongside renewables.")

# Image Right
img_path = r"C:\Users\aloks\.gemini\antigravity-ide\brain\ddd039d2-cd98-4d4a-876a-8a7511b866fd\fuel_sources_tab_1784127642558.png"
slide.shapes.add_picture(img_path, Inches(5.8), Inches(1.5), width=Inches(6.8))


# SLIDE 6: Dashboard Overview (Dark)
slide = prs.slides.add_slide(slide_layout)
set_bg_color(slide, BG_DARK)
add_slide_header(slide, "Interactive Dashboard Architecture", dark_mode=True)

# Image Left
img_path = r"C:\Users\aloks\.gemini\antigravity-ide\brain\ddd039d2-cd98-4d4a-876a-8a7511b866fd\global_trends_tab_1784127624483.png"
slide.shapes.add_picture(img_path, Inches(0.8), Inches(1.5), width=Inches(6.5))

# Text Right
right_box = slide.shapes.add_textbox(Inches(7.5), Inches(1.5), Inches(5.0), Inches(5.0))
tf = right_box.text_frame
tf.word_wrap = True

def add_bullet_light(tf, title, desc, space=12):
    p = tf.add_paragraph() if tf.paragraphs[0].text else tf.paragraphs[0]
    p.text = title + ": "
    p.font.bold = True
    p.font.size = Pt(15)
    p.font.color.rgb = TEXT_LIGHT
    p.font.name = 'Arial'
    
    run = p.add_run()
    run.text = desc
    run.font.bold = False
    run.font.size = Pt(15)
    run.font.color.rgb = MUTED_GRAY
    
    p.space_after = Pt(space)

add_bullet_light(tf, "Sleek Dark Theme", "Designed with a radial dot-grid pattern and responsive glassmorphic widgets to deliver a premium user experience.")
add_bullet_light(tf, "Dynamic Filters", "Users can toggle years, regions, and select specific countries to recalculate variables and index indices dynamically.")
add_bullet_light(tf, "Hybrid Layout", "Charts are nested inside high-contrast white container boxes to maintain maximum visual clarity and standard compliance.")


# SLIDE 7: Summary & Conclusions (Dark)
slide = prs.slides.add_slide(slide_layout)
set_bg_color(slide, BG_DARK)
add_slide_header(slide, "Conclusions & Key Takeaways", dark_mode=True)

box = slide.shapes.add_textbox(Inches(1.5), Inches(1.8), Inches(10.3), Inches(4.5))
tf = box.text_frame
tf.word_wrap = True

add_bullet_light(tf, "1. The Decoupling Proof", "Economic development no longer strictly requires increasing carbon emissions. Policies, renewable investments, and efficiency upgrades have successfully bent the curve in major developed nations.", space=20)
add_bullet_light(tf, "2. Asia's Central Role", "Global decarbonization cannot occur without support for industrial transition in Asia. Providing technology sharing and capital to help developing economies leapfrog coal is critical.", space=20)
add_bullet_light(tf, "3. Storytelling with Data", "Combining structured exploratory analysis (notebook) and dynamic, customer-facing delivery (dashboard) ensures insights are clear, compelling, and actionable for decision-makers.", space=20)


prs.save('presentation.pptx')
print("Successfully generated presentation.pptx")
